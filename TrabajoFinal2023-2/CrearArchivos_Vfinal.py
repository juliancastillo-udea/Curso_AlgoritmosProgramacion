# -*- coding: utf-8 -*-
"""
Created on Wed Aug 30 11:12:25 2023
Trabajo Final Algoritmia y Programación
@author: JulianCastillo
Meaning: El presente documento crea 1000 archivos cortos de multiples formatos
y los aloja en una carpeta especifica.
"""

import os
from openpyxl import Workbook
from docx import Document
import random as rnd
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import time
import datetime
import warnings
warnings.filterwarnings('ignore')
import logging
import tqdm

inicio = time.time()
print('Ingresar la cantidad de archivos a procesar, ingresar solo numeros enteros.')
print('Si ingresa el numero 10, por ejemplo, se crearan 10 archivos de cada uno.')
print('Generando un total de 100 archivos, A continuacion ingresar el numero.')
archivos = int(input('Ingresar la cantidad: '))
hoy = datetime.date.today().strftime('%Y%m%d')
nombre_archivo_log = f"log_{hoy}.log"
logging.basicConfig(filename=nombre_archivo_log, level=logging.INFO,
                    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logging.info("Iniciando el proceso, por CastilloEnterprises (^人^)(^人^)(^人^)")


#Creamos el directorio (carpeta) en donde se crearan los archivos
DirectorioActual = os.getcwd()
print(f'El directorio actual de trabajo es: \n\t--> {DirectorioActual}, \nEsta carpeta contendrá los archivos del trabajo final')
CarpetaNueva = "CarpetaArchivosTrabajoFinal"
os.makedirs(CarpetaNueva, exist_ok=True)
logging.info("Se crea el directorio {}".format(CarpetaNueva))
print(f"'{CarpetaNueva}' La carpeeta ha sido creada.")
# Nombre de la carpeta donde se crearán los archivos
carpeta = os.path.join(DirectorioActual, "CarpetaArchivosTrabajoFinal")
logging.info("La ruta de trabajo será {}".format(carpeta))
RutaNombres = r'NombresArgentina.csv'
RutaNombres = os.path.join(DirectorioActual, RutaNombres)
logging.info("Cargando CSV con nombres")
dfNombres = pd.read_csv(RutaNombres, encoding='ISO-8859-1')
Nombres = dfNombres['name'].tolist()
logging.info("Reemplazando nombres y detalles del documento")
for i in tqdm.trange(len(Nombres)):
    if ' ' in Nombres[i]:
        Nombres[i]=Nombres[i].replace(' ', '_')
logging.info("Creacion de archivos")
for i in tqdm.trange(int(archivos), miniters=int(1), ascii=True,desc="Creando Archivos", dynamic_ncols=True):
    # Crear un archivo de Excel usando openpyxl--------------------------------
    logging.info(f"Creacion de archivos-->Excel {i+1}")
    wb = Workbook()
    ws = wb.active
    ws['A1'] = 'Hola soy un numero aleatorio'
    ws['A2'] = rnd.random()
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    excel_filename = os.path.join(carpeta, tempNombre+'.xlsx')
    wb.save(excel_filename)
    
    # Crear un archivo de Word usando python-docx------------------------------
    logging.info(f"Creacion de archivos-->Word {i+1}")
    doc = Document()
    doc.add_paragraph(f"¡Hola! Este es un archivo de Word creado desde Python. \n Soy un numero aleatorio {rnd.random()}")
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    word_filename = os.path.join(carpeta, tempNombre+'.docx')
    doc.save(word_filename)
    
    # Crear un archivo de texto (txt)------------------------------------------
    logging.info(f"Creacion de archivos-->Texto {i+1}")
    texto = f"¡Hola! Este es un archivo de Word creado desde Python. \n Soy un numero aleatorio {rnd.random()}"
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    txt_filename = os.path.join(carpeta, tempNombre+'.txt')
    with open(txt_filename, 'w') as txt_file:
        txt_file.write(texto)
    
    # Crear un archivo CSV-----------------------------------------------------
    logging.info(f"Creacion de archivos-->CSV {i+1}")
    csv_data = f"Nombre,Edad\n{rnd.choice(Nombres)},{rnd.uniform(18,60)}\n"
    csv_data += f"{rnd.choice(Nombres)},{rnd.uniform(18,60)}\n"
    csv_data += f"{rnd.choice(Nombres)},{rnd.uniform(18,60)}\n"
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    csv_filename = os.path.join(carpeta, tempNombre+'.csv')
    with open(csv_filename, 'w') as csv_file:
        csv_file.write(csv_data)
    
    # Crear un archivo JSON----------------------------------------------------
    logging.info(f"Creacion de archivos-->JSON {i+1}")
    import json
    json_data = {"nombre": rnd.choice(Nombres), "edad": rnd.uniform(18,60)}
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    json_filename = os.path.join(carpeta, tempNombre+'.json')
    with open(json_filename, 'w') as json_file:
        json.dump(json_data, json_file, indent=4)
    
    # Crear un archivo XML-----------------------------------------------------
    logging.info(f"Creacion de archivos-->XML {i+1}")
    xml_data = f"""<?xml version="1.0" encoding="UTF-8"?>
    <datos>
        <persona>
            <nombre>{rnd.choice(Nombres)}</nombre>
            <edad>{rnd.uniform(18,60)}</edad>
        </persona>
        <persona>
            <nombre>{rnd.choice(Nombres)}</nombre>
            <edad>{rnd.uniform(18,60)}</edad>
        </persona>
    </datos>
    """
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    xml_filename = os.path.join(carpeta, tempNombre+'.xml')
    with open(xml_filename, 'w') as xml_file:
        xml_file.write(xml_data)
    
    # Crear un archivo de PowerPoint-------------------------------------------
    logging.info(f"Creacion de archivos-->PowerPoint {i+1}")
    from pptx import Presentation
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = f"Presentación creada desde Python,\n Soy un numero aleatorio {rnd.random()}"
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    ppt_filename = os.path.join(carpeta, tempNombre+'.pptx')
    ppt.save(ppt_filename)
    
    # Crear un archivo HTML----------------------------------------------------
    logging.info(f"Creacion de archivos-->HTML {i+1}")
    html_data = """<!DOCTYPE html>
    <html>
    <head>
        <title>Página HTML creada desde Python</title>
    </head>
    <body>
        <h1>Hola desde Python</h1>
        <h1>Hola soy un numero aleatorio {}</h1>
    </body>
    </html>
    """.format(rnd.random())
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    html_filename = os.path.join(carpeta, tempNombre+'.html')
    with open(html_filename, 'w') as html_file:
        html_file.write(html_data)
    
    #Crear una imagen vacía----------------------------------------------------
    logging.info(f"Creacion de archivos-->Imagen PNG {i+1}")
    ancho = 400
    alto = 400   
    imagen_vacia = Image.new('RGB', (ancho, alto), 'white')
    texto = f"Hola soy un numero aleatorio\n {rnd.random()}"
    fuente = ImageFont.truetype("arial.ttf", 20)
    dibujar = ImageDraw.Draw(imagen_vacia)
    bbox = dibujar.textbbox((0,0), texto, font=fuente)
    ancho_texto = bbox[2] - bbox[0]
    alto_texto = bbox[3] - bbox[1]
    x = (ancho - ancho_texto) / 2
    y = (alto - alto_texto) / 2
    dibujar.text((x, y), texto, font=fuente, fill="black")
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    imagen = os.path.join(carpeta, tempNombre+'.png')
    imagen_vacia.save(imagen)
    
    #Archivo PDF---------------------------------------------------------------
    logging.info(f"Creacion de archivos-->PDF {i+1}")
    tempNombre = rnd.choice(Nombres)
    Nombres.remove(tempNombre)
    pdfarchivo = os.path.join(carpeta, tempNombre+'.pdf')
    c = canvas.Canvas(pdfarchivo, pagesize=letter)
    width, height = letter
    texto = f"Hola soy un archivo PDF con un numero aleatorio \n {rnd.random()}"
    c.drawString(width/4 - 15, height/3, texto)  # Ajusta la posición según lo que necesites
    c.save()
logging.info("Proceso Terminado......")
logging.info("Gracias por usar Castillo Enterprises")
logging.shutdown()

fin = time.time()
tiempoejecucion = fin - inicio
minutos = int(tiempoejecucion // 60)
segundos = tiempoejecucion % 60
print(f"El script tomó {minutos} minutos y {segundos:.2f} segundos para ejecutarse.")
